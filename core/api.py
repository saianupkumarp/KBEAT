from flask import Blueprint, abort, jsonify, request, current_app, Response, send_file
from core import data, tasks
from datetime import datetime
import os
import StringIO, mimetypes
from werkzeug.datastructures import Headers

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches

from reportlab.platypus import BaseDocTemplate, Frame, Paragraph, NextPageTemplate, PageBreak, PageTemplate, Table, TableStyle, Image, Flowable
from reportlab.graphics.shapes import Drawing, _DrawingEditorMixin
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.textlabels import Label
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.legends import Legend
from reportlab.lib.units import inch
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.pdfgen.canvas import Canvas
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.lib.validators import Auto
from reportlab.lib.pagesizes import A4
from reportlab.lib.enums import TA_CENTER
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
#Rest API
rest_api = Blueprint('rest_api', __name__)

def _paginate(items, **options):
    offset = options.get('offset', 0)
    count = options.get('count', 0)
    items = items[offset:offset + count if count > 0 else None]

    return dict(objects=items, metadata=dict(count=len(items), offset=offset))

@rest_api.route('/config')
def get_config():
    return jsonify(app_name=current_app.config['APP_NAME'], locales=data.get_locales())

@rest_api.route('/models')
def get_models():
    i18n = request.args.get('locale', 'en')
    return jsonify(_paginate([models.to_primitive() for models in data.get_models()]))

@rest_api.route('/model/<model_name>')
def get_model(model_name):
    model = data.get_model(model_name)
    i18n = request.args.get('locale', 'en')
    return jsonify(model.to_primitive()) if model else abort(404)

@rest_api.route('/models/<model_name>', methods=['POST'])
def run_model(model_name):
    return tasks.run_task(request.data, model_name)

@rest_api.route('/tasks')
def get_tasks(count=10, offset=0):
    return jsonify(data.get_tasks())

@rest_api.route('/tasks/<task_id>')
def get_task(task_id):
    task = data.get_task(task_id)
    return jsonify(task.to_primitive(role='DTO')) if task else abort(404)

@rest_api.route('/tasks/<task_id>/result')
def get_task_result(task_id):
    result = data.get_task_result(task_id)
    return jsonify(result) if result else abort(404)


@rest_api.route('/tasks/<task_id>/pptexport', methods=['GET', 'POST'])
def get_pptx_results(task_id):
    # Flask response
    response = Response()
    response.status_code = 200


    task = data.get_task_result(task_id)

    ## Create presentation with given template
    prs = Presentation("templates/ppt_template.pptx")
    # prs = Presentation()
    # slide_master = prs.slide_master

    ## Define slide template /prs.slide_layouts[number of slide in slide matser template]/
    title_slide_layout = prs.slide_layouts[0]
    table_slide_layout = prs.slide_layouts[4]
    chart_slide_layout = prs.slide_layouts[4]

    ## Slide 1
    firstSlide = prs.slides.add_slide(title_slide_layout)
    ## Create shapes and add title
    firstSlideShapes = firstSlide.shapes
    firstSlideShapes.title.text = "KAPSARC Building Stock Energy Efficiency Analysis"

    ## Slide 2
    PSEtableSlide = prs.slides.add_slide(table_slide_layout)
    PSEtableSlideShapes = PSEtableSlide.shapes
    PSEtableSlideShapes.title.text = 'PSE Table'
    ## Get barChartData
    barChartData = task['barChartData']
    month=['JAN', 'FEB', 'MAR', 'APR', 'MAY', 'JUN', 'JUL', 'AUG', 'SEP', 'OCT', 'NOV', 'DEC']

    cols = len(month)+1
    rows = len(barChartData)+1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(0.8)
    ## Add table
    table = PSEtableSlideShapes.add_table(rows, cols, left, top, width, height).table
    # write column headings
    table.cell(0, 0).text = 'Key'
    table.columns[0].width = Inches(1.5)

    ## Add table header and apply text style
    for i,m in enumerate(month):
        table.cell(0, i+1).text = str(month[i])
        table.cell(0,i+1).text_frame.paragraphs[0].font.size=Pt(14)

    ## Create 'barData' list to pass the data values to the bar chart
    barData=[[0 for i in xrange(len(barChartData[0]['values']))] for i in xrange(len(barChartData))]
    ## Add table index and values
    for i,result in enumerate(barChartData):
        table.cell(i+1, 0).text = str(result['key'])
        table.cell(i+1,0).text_frame.paragraphs[0].font.size=Pt(14)
        for j,value in enumerate(result['values']):
            table.cell(i+1, j+1).text = str(result['values'][j]['y'])
            barData[i][j]=str(result['values'][j]['y'])
            table.cell(i+1,j+1).text_frame.paragraphs[0].font.size=Pt(12)



    ## Slide 3
    PSEchartSlide = prs.slides.add_slide(chart_slide_layout)
    PSEchartSlideShapes = PSEchartSlide.shapes
    PSEchartSlideShapes.title.text = 'PSE Chart'
    ## Define chart data
    chart_data = ChartData()
    chart_data.categories = month
    for i,result in enumerate(barChartData):
        chart_data.add_series(str(result['key']), barData[i])

    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    PSEchart=PSEchartSlideShapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, cx, cy, chart_data).chart
    # Add legend
    PSEchart.has_legend = True
    PSEchart.legend.position = XL_LEGEND_POSITION.RIGHT
    PSEchart.legend.include_in_layout = False


    ## Slide 4
    BEPUtableSlide = prs.slides.add_slide(table_slide_layout)
    BEPUtableSlideShapes = BEPUtableSlide.shapes
    BEPUtableSlideShapes.title.text = 'BEPU Table'
    pieChartData = task['pieChartData']

    cols = 2
    rows = len(pieChartData)+1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(5.5)
    height = Inches(1.5)

    table = BEPUtableSlideShapes.add_table(rows, cols, left, top, width, height).table

    # write column headings
    table.cell(0, 0).text = 'Key'
    table.cell(0, 1).text = 'Value'
    ## Create label and values lists to pass the data to the pie chart
    pieDataLabel=[0 for x in range(len(pieChartData))]
    pieDataValue=[0 for x in range(len(pieChartData))]

    ## Add table index and values
    for i,result in enumerate(pieChartData):
        # write body cells
        table.cell(i+1, 0).text = str(result['label'])
        table.cell(i+1, 1).text = str(result['value'])
        pieDataValue[i]=result['value']
        pieDataLabel[i]=str(result['label'])


    ## Slide 5
    chart_slide_layout = prs.slide_layouts[4]
    BEPUchartSlide = prs.slides.add_slide(chart_slide_layout)
    BEPUchartSlideShapes = BEPUchartSlide.shapes
    BEPUchartSlideShapes.title.text = 'BEPU Chart'

    # Define chart data
    chart_data = ChartData()
    chart_data.categories = pieDataLabel
    chart_data.add_series('', pieDataValue)
    # add chart to slide
    x, y, cx, cy = Inches(0.5), Inches(1), Inches(6), Inches(5.5)
    BEPUchart=BEPUchartSlideShapes.add_chart(XL_CHART_TYPE.PIE , x, y, cx, cy, chart_data).chart
    # Add legend
    BEPUchart.has_legend = True
    BEPUchart.legend.position = XL_LEGEND_POSITION.RIGHT
    BEPUchart.legend.include_in_layout = False
    BEPUchart.plots[0].has_data_labels = True
    data_labels = BEPUchart.plots[0].data_labels
    # data_labels=labelc
    # data_labels.number_format_is_linked = False
    # data_labels.number_format = '0%' #%
    data_labels.ShowPercentage=True
    data_labels.ShowValue=False
    data_labels.position = XL_LABEL_POSITION.CENTER



    ## Slide 6
    title_only_slide_layout = prs.slide_layouts[4]
    CalibrationSlide = prs.slides.add_slide(title_only_slide_layout)
    CalibrationSlideShapes = CalibrationSlide.shapes
    CalibrationSlideShapes.title.text = 'Calibration Data'
    #get data values
    calibrationData = task['calibrationData']

    cols = 2
    rows = len(calibrationData)+1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(7.5)
    height = Inches(0.8)

    table = CalibrationSlideShapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(2.3)
    calibrationDataLabel=[0 for x in range(len(calibrationData))]
    for i,result in enumerate(calibrationData):
        # write body cells
        table.cell(i, 0).text = str(result)
        table.cell(i, 1).text = str(calibrationData[result])


    #Saving file to a in-memory file
    output_file = StringIO.StringIO()
    prs.save(output_file)
    output_file.seek(0)

    # Set filname and mimetype
    file_name = 'K-BEAT_export_{}.pptx'.format(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

    #Returning the file from memory
    return send_file(output_file, attachment_filename=file_name, as_attachment=True)




@rest_api.route('/tasks/<task_id>/pdfexport', methods=['GET', 'POST'])
def get_pdf_results(task_id):

        # Flask response
    response = Response()
    response.status_code = 200


    task = data.get_task_result(task_id)
    #Saving file to a in-memory file
    output_file = StringIO.StringIO()

    def header_footer(canvas, doc):

        canvas.saveState()

        background ='static/img/pdf_bg.png'
        canvas.drawImage(background,1*inch,5.75*inch, width=8*inch,height=6*inch,mask='auto')


        # Header
        logo =Image('static/img/logo/logo.png')
        logo.drawHeight = 0.5*inch
        logo.drawWidth = 1.75*inch
        date=datetime.now().strftime("%y-%m-%d %H:%M")
        headerData= [[logo, '', date]]
        headerTable = Table(headerData, colWidths=[2*inch,3.58 * inch,1.2* inch],
             style=[('LINEBELOW',(0,0),(2,0),1,colors.HexColor(0xcccccc)),
                    ('TEXTCOLOR',(0,0),(2,0),colors.HexColor(0x807F83)),
                    ('VALIGN',(1,0),(1,0),'MIDDLE'),
                    ('VALIGN',(2,0),(2,0),'MIDDLE')])
        headerTable.wrapOn(canvas, doc.width, doc.topMargin)
        headerTable.drawOn(canvas, doc.leftMargin, doc.height + doc.topMargin)

        pageNum ="Page %d" % doc.page
        footerData= [['KAPSARC Building Energy Assessment Tool (BEAT)', pageNum]]
        footerTable = Table(footerData, colWidths=[5.76*inch,1* inch],
             style=[('LINEABOVE',(0,0),(1,0),2,colors.HexColor(0xcccccc)),
                    ('TEXTCOLOR',(0,0),(1,0),colors.HexColor(0x807F83)),
                    ('ALIGN',(1,0),(1,0),'RIGHT')])
        footerTable.wrapOn(canvas, doc.width, doc.bottomMargin)
        footerTable.drawOn(canvas, doc.leftMargin, 0.5*inch)

        canvas.restoreState()

    pdfmetrics.registerFont(TTFont('Vera', 'Vera.ttf'))
    pdfmetrics.registerFont(TTFont('VeraBd', 'VeraBd.ttf'))
    pdfmetrics.registerFont(TTFont('VeraIt', 'VeraIt.ttf'))
    pdfmetrics.registerFont(TTFont('VeraBI', 'VeraBI.ttf'))
    styles=getSampleStyleSheet()
    # Title
    styles.add(ParagraphStyle(name = 'styleTitle',
                                      alignment= TA_CENTER,
                                      fontSize = 16,
                                      fontName ='Vera',
                                      textColor= colors.HexColor(0x61a659),
                                      leading = 30,
                                      spaceBefore = 35,
                                      spaceAfter = 10))
    # Headings
    styles.add(ParagraphStyle(name = 'styleHeading',
                                      parent = styles['Heading2'],
                                      fontSize = 14,
                                      textColor = colors.HexColor(0x807F83),
                                      leading = 20,
                                      spaceBefore = 10,
                                      underlineProportion=1.1,
                                      spaceAfter = 10))
    styles.add(ParagraphStyle(name = 'styleHeading2',
                                      parent = styles['Heading2'],
                                      fontSize = 14,
                                      textColor= colors.HexColor(0x61a659),
                                      leading = 20,
                                      spaceBefore = 20,
                                      underlineProportion=1.1,
                                      spaceAfter = 20))
    styles.add(ParagraphStyle(name = 'styleHeading3',
                                      #alignment= TA_CENTER,
                                      fontSize = 12,
                                      fontName ='Vera',
                                      textColor= colors.HexColor(0x61a659),
                                      leading = 20,
                                      spaceBefore = 10,
                                      spaceAfter = 5))
    # Body text
    styles.add(ParagraphStyle(name = 'styleBodyText',
                                      parent = styles['Normal'],
                                      fontSize = 9,
                                      textColor= colors.HexColor(0x666666),
                                      spaceBefore = 5,
                                      spaceAfter = 15))

    styleTitle = styles['styleTitle']
    styleHeading = styles['styleHeading']
    styleHeading2 = styles['styleHeading2']
    styleHeading3 = styles['styleHeading3']
    styleBodyText = styles['styleBodyText']
    pdf_chart_colors = [
            "#3D6531",
            "#61a24f",
            "#89B97B",
            "#B0D1A7",
            "#cde5c7",
            "#7e7f82",
            "#9E9FA1",
            "#BFBFC1",
            "#DFDFE0",
            "#ffd200",
            "#FFE360",
            "#FFEE9F",
            ]

    Elements=[]
    doc = BaseDocTemplate(output_file, showBoundary=0, pagesize=A4,title='KASPSARC BEAT Report',author="KAPSARC",
    leftMargin=0.75*inch, rightMargin=0.75*inch, topMargin=inch, bottomMargin=inch)

    frame = Frame(doc.leftMargin, doc.topMargin, doc.width, doc.height,topPadding=0.3*inch, showBoundary=0)
    template = PageTemplate(id='template', frames=[frame], onPage=header_footer)
    doc.addPageTemplates([template])


    ## PAGE 1
    #add some flowables
    Elements.append(Paragraph("KAPSARC Building Energy Assessment Tool (BEAT)",styleTitle))
    Elements.append(Paragraph("Your Building Description",styleHeading))
    rowHeights=0.3*inch
    calibrationData=task['calibrationData']
    Elements.append(Paragraph("General Information:",styleHeading3))
    infoTableData=[[Paragraph('<b>- Name: </b>'+calibrationData['txtBldgName'],styleBodyText),
                   Paragraph('<b>- Address: </b>'+calibrationData['txtBldgAddress'],styleBodyText),
                   Paragraph('<b>- Type: </b>'+calibrationData['cmbBldgType'],styleBodyText)],
                   [Paragraph('<b>- Location: </b>'+calibrationData['cmbBldgLocation'],styleBodyText),
                   Paragraph('<b>- Shape: </b>'+calibrationData['cmbBldgShape'],styleBodyText),
                   Paragraph('<b>- Floor Area (m'+u"\u00b2"+'): </b>'+str(calibrationData['txtFloorArea']),styleBodyText)]
    ]
    infoTable=Table(infoTableData, colWidths=[160,165,150], rowHeights=rowHeights)
    Elements.append(infoTable)
  
    Elements.append(Paragraph('<br />', styleBodyText)) 
    Elements.append(Paragraph("Envelope Construction Details:",styleHeading3))
    envTableData=[[Paragraph('<b>- South Wall: </b>'+calibrationData['cmbSouthWall'],styleBodyText),
                  Paragraph('<b>- West Wall: </b>'+calibrationData['cmbWestWall'],styleBodyText)],
                  [Paragraph('<b>- North Wall: </b>'+calibrationData['cmbNorthWall'],styleBodyText),
                  Paragraph('<b>- East Wall: </b>'+calibrationData['cmbEastWall'],styleBodyText)],
                  [Paragraph('<b>- Roof: </b>'+calibrationData['cmbRoof'],styleBodyText),
                  Paragraph('<b>- Floor: </b>'+calibrationData['cmbFirstFloorContact'],styleBodyText)],
                  [Paragraph('<b>- Windows Type: </b>'+calibrationData['glasstype'],styleBodyText),
                  Paragraph('<b>- Overhang Depth (m): </b>'+str(calibrationData['txtWinSouthOverhang']),styleBodyText)]
    ]
    envTable=Table(envTableData,colWidths=[240,235], rowHeights=rowHeights)
    Elements.append(envTable)

    Elements.append(Paragraph('<br />', styleBodyText)) 
    Elements.append(Paragraph("Air Conditioning Systems",styleHeading3))
    hvacTableData=[[Paragraph('<b>- HVAC  System Type: </b>'+calibrationData['cmbBldgSystem'],styleBodyText),
                    Paragraph('<b>- Cooling Temperature Setting ('+u"\u00b0"+'C): </b>'+str(calibrationData['txtCoolSetTemp']),styleBodyText)],
                    [Paragraph('<b>- Energy Efficiency Ratio (EER): </b>'+str(calibrationData['eir']),styleBodyText),
                    Paragraph('<b>- Heating Temperature Setting ('+u"\u00b0"+'C): </b>'+str(calibrationData['txtHeatSetTemp']),styleBodyText)]
    ]
    hvacTable=Table(hvacTableData,colWidths=[240,235],rowHeights=rowHeights)
    Elements.append(hvacTable)

    Elements.append(Paragraph('<br />', styleBodyText)) 
    Elements.append(Paragraph("Overall Assessment",styleHeading))
    Elements.append(Paragraph("Based on your description and current SASO requirements, the tool provides the following assessments:",styleBodyText))    
    
    if task['compliant']:
        compliant="<strong><font color=green>meets</font></strong>"
    else: 
        compliant="<strong><font color=red>does not meet</font></strong>"
    if (task['ngEnergyDiff']<0):
        energyDiff= "<strong><font color=green>"+str(task['energyDiff'])+" kWh/year, less</font></strong>"
    else:
        energyDiff= "<strong><font color=red>"+str(task['energyDiff'])+" kWh/year, more</font></strong>"
    
    Elements.append(Paragraph("- Your building "+compliant+" SASO requirements for all building envelope",styleBodyText)) 
    Elements.append(Paragraph("- Your building consumed " +energyDiff+" than the SASO Baseline",styleBodyText)) 

    if task['compliant']:
        Elements.append(Paragraph("- You may reduce even more your energy consumption in your building by using LED lamps and high efficient appliances and air conditioning system",styleBodyText))  
    else:
        Elements.append(Paragraph(" - You need to add more insulation to the walls and/or roof, or use more efficient window glazing to comply with SASO requirements",styleBodyText))  
    if  not task['compliant'] and (task['ngEnergyDiff']>=0):
        Elements.append(Paragraph(" - You may also consider using LED lamps and energy efficient appliances and air conditioning system",styleBodyText))   
    

    Elements.append(PageBreak())
    Elements.append(Paragraph("How electricity is used in your building?",styleHeading3))
    Elements.append(Paragraph("Your building needs electricity to operate several types of equipment including: air-conditioning, lighting, appliances and domestic hot water.",styleBodyText))
    
    
    #add image
    Elements.append(Image('static/img/results-intro.png', width=4*inch, height=1.2*inch))
    #add text
    Elements.append(Paragraph("Based on the description you provided as well as typical lighting and appliances used in households, here how your building consumes electricity on annual basis:",styleBodyText))
    
    
    bepuPieData = task['bepuPieData']
    bepuTableData=[[0 for i in xrange(len(bepuPieData[0]))] for i in xrange(len(bepuPieData)+1)]
    bepuChartLabel=[0 for i in xrange(len(bepuPieData))]
    bepuChartData=[0 for i in xrange(len(bepuPieData))]
    bepuTableData[0][0]= 'End-Use'
    bepuTableData[0][1] =  'Annual Electricity Use'
    for i,result in enumerate(bepuPieData):
        # write body cells
        bepuTableData[i+1][0]= str(result['label'])
        bepuTableData[i+1][1] = int(result['value'])
        bepuChartLabel[i] = str(result['label'])
        bepuChartData[i] = result['value']

    #add chart
    bepuChart = Drawing(400,200)
    pc = Pie()
    pc.data = bepuChartData
    labelc=[0 for i in xrange(len(bepuChartData))]
    for i , r in enumerate(bepuChartData):
        labelc[i]=str(round(r/sum(bepuChartData)*100,1))+"%"
    pc.labels = labelc
    pc._seriesCount = len(bepuChartLabel)
    pc.slices.strokeColor     = colors.HexColor(0xffffff)
    pc.slices.strokeWidth     = 0.5
    bepu_chart_colors=['#FFC43E','#A4A4A4','#F67A40','#5894D0','#98cd99']
    for i , r in enumerate(bepuChartLabel):
        pc.slices[i].fillColor= colors.HexColor(bepu_chart_colors[i])

    pc.width=pc.height= 120
    pc.x=40
    pc.y=30
    # add_legend(d, pc)
    legend = Legend()
    legend.alignment = 'right'
    legend.x = pc.width+pc.x+80
    legend.y = pc.height-10
    legend.dx              = 8
    legend.dy              = 8
    legend.fontName        = 'Helvetica'
    legend.fillColor = colors.HexColor(0x807F83)
    legend.fontSize        = 10
    legend.boxAnchor       = 'nw'
    legend.columnMaximum   = 8
    legend.strokeWidth     = 0.5
    legend.strokeColor     = colors.HexColor(0xffffff)
    legend.deltax          = 75
    legend.deltay          = 10
    legend.autoXPadding    = 10
    legend.yGap            = 0
    legend.dxTextSpace     = 5
    legend.dividerLines    = 1|2|4
    legend.dividerOffsY    = 6
    legend.subCols.rpad    = 70
    legend.dividerColor    = colors.HexColor(0xdedede)
    legend.colorNamePairs = [(pc.slices[i].fillColor, (bepuChartLabel[i][0:20], '  %s ' % "{:,}".format(int(pc.data[i])))) for i in xrange(len(pc.data))]
    legendHeader = Legend()
    legendHeader.colorNamePairs = [('', ('End-Use','Annual Electricity Use\n(kWh/year)'))]
    legendHeader.alignment = 'right'
    legendHeader.x = legend.x-20
    legendHeader.y = legend.y+30
    legendHeader.fontName        = 'Helvetica'
    legendHeader.fillColor = colors.HexColor(0x807F83)
    legendHeader.fontSize = 10
    legendHeader.boxAnchor       = 'nw'
    legendHeader.subCols.rpad= 80
    legendFooter = Legend()
    legendFooter.colorNamePairs = [('', ('Total',str("{:,}".format(int(sum(bepuChartData))))+''))]
    legendFooter.alignment = 'right'
    legendFooter.x = legendHeader.x+5
    legendFooter.y = legend.y-(len(bepuChartLabel)+1)*10
    legendFooter.fontName        = 'Helvetica-Bold'
    legendFooter.fillColor = colors.HexColor(0x807F83)
    legendFooter.fontSize = 10
    legendFooter.boxAnchor       = 'nw'
    legendFooter.subCols.rpad    = 145
    bepuChart.add(legend)
    bepuChart.add(legendHeader)
    bepuChart.add(legendFooter)
    pc.slices.fontColor = colors.HexColor(0x807F83)
    n = len(pc.data)
    bepuChart.add(pc, '')
    Elements.append(bepuChart)

    ## PAGE 2
    Elements.append(Paragraph("When electricity is consumed in your building?",styleHeading3))
    Elements.append(Paragraph("Based on the weather of your location as well as typical lighting and appliances used in households, your building consumes electricity as noted in the following monthly profile:",styleBodyText))    
    #add chart
    pseBarData = task['pseBarData']

    pseTableData=[[0 for i in xrange(len(pseBarData[0]['values'])+1)] for i in xrange(len(pseBarData)+1)]
    pseChartData=[[0 for i in xrange(len(pseBarData[0]['values']))] for i in xrange(len(pseBarData))]
    pseChartLegend=[0 for i in xrange(len(pseBarData))]
    pseTableData[0][0] = 'Key'
    month=['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    for i,m in enumerate(month):
        pseTableData[0][i+1] = str(month[i])
    for i,result in enumerate(pseBarData):
        # write body cells
        pseTableData[i+1][0] = str(result['key'])
        pseChartLegend[i] = str(result['key'])
        for j,value in enumerate(result['values']):
            pseTableData[i+1][j+1]=int(result['values'][j]['y'])
            pseChartData[i][j]=int(result['values'][j]['y'])

    pseChart = Drawing(400, 200)
    bc = VerticalBarChart()
    bc.x=70
    bc.y=0
    bc.height = 200
    bc.width = 300
    bc.data = pseChartData
    bc.strokeColor = colors.black
    bc.valueAxis.valueMin = 0
    bc.strokeWidth = 0
    bc.valueAxis.valueMin = 0
    bc.categoryAxis.style = 'stacked'
    bc.categoryAxis.labels.boxAnchor = 'ne'
    bc.categoryAxis.labels.dx = 10
    bc.categoryAxis.labels.dy = -2
    bc.valueAxis.labels.fontName  = 'Helvetica'
    bc.valueAxis.labels.fontSize  = 10
    bc.valueAxis.strokeWidth  = 0.5
    bc.valueAxis.strokeColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.strokeWidth  = 0.5
    bc.categoryAxis.strokeColor  = colors.HexColor(0x807F83)
    bc.valueAxis.labels.fillColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.labels.fontName  = 'Helvetica'
    bc.categoryAxis.labels.fontSize  = 10
    bc.categoryAxis.labels.fillColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.categoryNames = month
    # create a list and add the elements of our document (image, paragraphs, table, chart) to it
    #add our barchart and legend
    bc.barWidth = .3*inch
    bc.groupSpacing = .2 * inch

    bc.bars.strokeColor     = colors.HexColor(0xffffff)
    bc.bars.strokeWidth     = 0.
    pse_chart_colors=['#FFC43E','#A4A4A4','#F67A40','#5894D0']
    for i , r in enumerate(pseChartLegend):
        bc.bars[i].fillColor= colors.HexColor(pse_chart_colors[i])
    #  = colors.blue
    legend = Legend()
    legend.alignment = 'right'
    legend.x = bc.width+bc.x+5
    legend.y = bc.height+bc.y
    legend.deltax = 40
    legend.dxTextSpace = 5
    legend.dx              = 8
    legend.dy              = 8
    legend.fontName        = 'Helvetica'
    legend.fillColor = colors.HexColor(0x807F83)
    legend.fontSize        = 10
    legend.boxAnchor       = 'nw'
    legend.columnMaximum   = (len(bc.data)+1)/2
    legend.strokeWidth     = 0.5
    legend.strokeColor     = colors.HexColor(0xffffff)
    legend.deltax          = 75
    legend.deltay          = 12
    legend.dividerColor    = colors.HexColor(0xdedede)
    legend.columnMaximum = len(pseChartLegend)
    legend.colorNamePairs = [(bc.bars[i].fillColor, pseChartLegend[i]) for i in xrange(len(bc.data))]
    #pseChart.hAlign = 'RIGHT'

    label = Label()
    label.setOrigin(10,bc.height/2)
    #label.boxAnchor = 'sw'
    label.angle = 90
    label.fillColor = colors.HexColor(0x807F83)
    label.setText('Electricity Consumption (kWh)')
    label.fontName        = 'Helvetica'

    pseChart.add(legend, 'legend')
    pseChart.add(bc)
    pseChart.add(label)    
    Elements.append(pseChart)

    Elements.append(PageBreak())

    ## PAGE 3
    Elements.append(Paragraph("Does your building meet SASO Thermal Performance Requirements?",styleHeading3))
    Elements.append(Paragraph("Based on your description, the thermal transmittance properties of the walls, roof and glazing are calculated, and compared with SASO thermal building performance requirements:",styleBodyText))    
          
    #add chart
    lvdData= task['lvdData']
    lvdChartData=[[0 for i in xrange(len(lvdData[0]['values']))] for i in xrange(len(lvdData))]
    lvdChartCategoryNames=[0 for i in xrange(len(lvdData[0]['values']))]
    lvdComparedObjKey=[0 for i in xrange(len(lvdData))]
    for i,result in enumerate(lvdData):
        # write body cells
        lvdComparedObjKey[i]= str(lvdData[i]['key'])
        for j,value in enumerate(result['values']):
            lvdChartCategoryNames[j]=value['label']
            lvdChartData[i][j]=value['value']

    

    lvdChart = Drawing(400, 200)
    bc = VerticalBarChart()
    bc.x=70
    bc.y=0
    bc.height = 200
    bc.width = 300
    bc.data = lvdChartData
    bc.strokeColor = colors.black
    # bc.fillColor=colors.blue
    bc.valueAxis.valueMin = 0
    bc.strokeWidth = 0
    bc.valueAxis.valueMin = 0
    bc.categoryAxis.labels.boxAnchor = 'n'
    bc.categoryAxis.labels.dx = 0
    bc.categoryAxis.labels.dy = -2
    # bc.categoryAxis.labels.angle = 20
    bc.valueAxis.labels.fontName  = 'Helvetica'
    bc.valueAxis.labels.fontSize  = 10
    bc.valueAxis.strokeWidth  = 0.5
    bc.valueAxis.strokeColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.strokeWidth  = 0.5
    bc.categoryAxis.strokeColor  = colors.HexColor(0x807F83)
    bc.valueAxis.labels.fillColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.labels.fontName  = 'Helvetica'
    bc.categoryAxis.labels.fontSize  = 8
    bc.categoryAxis.labels.fillColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.categoryNames = lvdChartCategoryNames
    bc.categoryAxis.labels.angle= 0
    # create a list and add the elements of our document (image, paragraphs, table, chart) to it
    #add our barchart and legend
    bc.barWidth = .3*inch
    bc.groupSpacing = .2 * inch

    bc.bars.strokeColor     = colors.HexColor(0xffffff)
    bc.bars.strokeWidth     = 0.5
    lvd_chart_colors=['#5894D0','#F67A40']
    for i , r in enumerate(lvdComparedObjKey):
        bc.bars[i].fillColor= colors.HexColor(lvd_chart_colors[i])
    
    #  = colors.blue
    legend = Legend()
    legend.alignment = 'right'
    legend.x = bc.width+bc.x+5
    legend.y = bc.height+bc.y
    legend.deltax = 40
    legend.dxTextSpace = 5
    legend.dx              = 8
    legend.dy              = 8
    legend.fontName        = 'Helvetica'
    legend.fillColor = colors.HexColor(0x807F83)
    legend.fontSize        = 10
    legend.boxAnchor       = 'nw'
    legend.columnMaximum   = (len(bc.data)+1)/2
    legend.strokeWidth     = 0.5
    legend.strokeColor     = colors.HexColor(0xffffff)
    legend.deltax          = 75
    legend.deltay          = 12
    legend.dividerColor    = colors.HexColor(0xdedede)
    legend.columnMaximum = len(lvdComparedObjKey)
    legend.colorNamePairs = [(bc.bars[i].fillColor, lvdComparedObjKey[i]) for i in xrange(len(bc.data))]
    #pseChart.hAlign = 'RIGHT'

    
    label = Label()
    label.setOrigin(10,bc.height/2)
    #label.boxAnchor = 'sw'
    label.angle = 90
    label.fillColor = colors.HexColor(0x807F83)
    label.setText('Envelope U-value (W/m'+u'\u00b2'+'.k)')
    label.fontName        = 'Helvetica'

    lvdChart.add(label)  
    lvdChart.add(legend, 'legend')
    lvdChart.add(bc)
    Elements.append(lvdChart)

    #Elements.append(PageBreak())
    Elements.append(Paragraph('<br /><br />', styleBodyText))   
    ## PAGE 4
    Elements.append(Paragraph("How energy efficient is your building?",styleHeading3))
    Elements.append(Paragraph("Using your input specifications, the annual electricity consumption is calculated and compared with a similar building that meets SASO requirements:",styleBodyText))    
    
    #add chart
    bepuComparisonData= task['bepuComparisonData']

    bepuComparisonChartData=[[0 for i in xrange(len(bepuComparisonData[0]['values']))] for i in xrange(len(bepuComparisonData))]
    bepuChartCategoryNames=[0 for i in xrange(len(bepuComparisonData[0]['values']))]
    bepuComparedObjKey=[0 for i in xrange(len(bepuComparisonData))]
    for i,result in enumerate(bepuComparisonData):
        # write body cells
        bepuComparedObjKey[i]= str(bepuComparisonData[i]['key'])
        for j,value in enumerate(result['values']):
            bepuChartCategoryNames[j]=value['label']
            bepuComparisonChartData[i][j]=value['value']

    

    bepuComparisonChart = Drawing(400, 200)
    bc = VerticalBarChart()
    bc.x=70
    bc.y=0
    bc.height = 200
    bc.width = 300
    bc.data = bepuComparisonChartData
    bc.strokeColor = colors.black
    # bc.fillColor=colors.blue
    bc.valueAxis.valueMin = 0
    bc.strokeWidth = 0
    bc.valueAxis.valueMin = 0
    bc.categoryAxis.labels.boxAnchor = 'ne'
    bc.categoryAxis.labels.dx = 10
    bc.categoryAxis.labels.dy = -2
    # bc.categoryAxis.labels.angle = 20
    bc.valueAxis.labels.fontName  = 'Helvetica'
    bc.valueAxis.labels.fontSize  = 10
    bc.valueAxis.strokeWidth  = 0.5
    bc.valueAxis.strokeColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.strokeWidth  = 0.5
    bc.categoryAxis.strokeColor  = colors.HexColor(0x807F83)
    bc.valueAxis.labels.fillColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.labels.fontName  = 'Helvetica'
    bc.categoryAxis.labels.fontSize  = 10
    bc.categoryAxis.labels.fillColor  = colors.HexColor(0x807F83)
    bc.categoryAxis.categoryNames = bepuChartCategoryNames
    bc.categoryAxis.labels.angle= 30
    # create a list and add the elements of our document (image, paragraphs, table, chart) to it
    #add our barchart and legend
    bc.barWidth = .3*inch
    bc.groupSpacing = .2 * inch

    bc.bars.strokeColor     = colors.HexColor(0xffffff)
    bc.bars.strokeWidth     = 0.5
    bepu_chart_colors=['#5894D0','#F67A40']
    for i , r in enumerate(bepuComparedObjKey):
        bc.bars[i].fillColor= colors.HexColor(bepu_chart_colors[i])
    #  = colors.blue
    # bc.bars[1].fillColor = colors.lightblue
    legend = Legend()
    legend.alignment = 'right'
    legend.x = bc.width+bc.x+5
    legend.y = bc.height+bc.y
    legend.deltax = 40
    legend.dxTextSpace = 5
    legend.dx              = 8
    legend.dy              = 8
    legend.fontName        = 'Helvetica'
    legend.fillColor = colors.HexColor(0x807F83)
    legend.fontSize        = 10
    legend.boxAnchor       = 'nw'
    legend.columnMaximum   = (len(bc.data)+1)/2
    legend.strokeWidth     = 0.5
    legend.strokeColor     = colors.HexColor(0xffffff)
    legend.deltax          = 75
    legend.deltay          = 12
    legend.dividerColor    = colors.HexColor(0xdedede)
    legend.columnMaximum = len(bepuComparedObjKey)
    legend.colorNamePairs = [(bc.bars[i].fillColor, bepuComparedObjKey[i]) for i in xrange(len(bc.data))]
    #pseChart.hAlign = 'RIGHT'

    label = Label()
    label.setOrigin(10,bc.height/2)
    #label.boxAnchor = 'sw'
    label.angle = 90
    label.fillColor = colors.HexColor(0x807F83)
    label.setText('Annual Energy Use (kWh/year)')
    label.fontName        = 'Helvetica'

    bepuComparisonChart.add(label)  
    bepuComparisonChart.add(legend, 'legend')
    bepuComparisonChart.add(bc)
    Elements.append(bepuComparisonChart)
    
    Elements.append(PageBreak())
    
    doc.build(Elements)

    output_file.seek(0)

    # Set filname and mimetype
    file_name = 'K-BEAT_export_{}.pdf'.format(datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

    #Returning the file from memory
    return send_file(output_file, attachment_filename=file_name, as_attachment=True)
